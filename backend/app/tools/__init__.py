from typing import Dict, Any
import json
import ast
import math
import re
import httpx
import os
import uuid
from datetime import datetime, timedelta
from urllib.parse import urlparse, parse_qs


def tool_calculator(args: Dict[str, Any]) -> Dict[str, Any]:
    """Safe calculator using AST parsing."""
    expr = str(args.get("expression", ""))
    try:
        tree = ast.parse(expr, mode='eval')
        for node in ast.walk(tree):
            if isinstance(node, (ast.Expression, ast.BinOp, ast.UnaryOp, ast.Constant,
                                 ast.Add, ast.Sub, ast.Mult, ast.Div, ast.Pow, ast.Mod,
                                 ast.FloorDiv, ast.USub, ast.UAdd)):
                continue
            raise ValueError(f"Unsupported operation: {type(node).__name__}")
        result = eval(compile(tree, '<calc>', 'eval'), {"__builtins__": {}})  # noqa: S307
        return {"result": result}
    except Exception as e:
        return {"error": str(e)}


def tool_youtube_transcript(args: Dict[str, Any]) -> Dict[str, Any]:
    """Extract transcript from a YouTube video — no API key needed."""
    from youtube_transcript_api import YouTubeTranscriptApi
    url = str(args.get("url", ""))
    try:
        # Parse video ID from various YouTube URL formats
        video_id = None
        parsed = urlparse(url)
        if "youtu.be" in (parsed.hostname or ""):
            video_id = parsed.path.lstrip("/").split("?")[0]
        elif "youtube.com" in (parsed.hostname or ""):
            qs = parse_qs(parsed.query)
            video_id = qs.get("v", [None])[0]
        # Handle bare video IDs
        if not video_id and re.match(r'^[a-zA-Z0-9_-]{11}$', url):
            video_id = url
        if not video_id:
            return {"error": f"Could not extract video ID from: {url}"}

        # Fetch transcript using youtube-transcript-api
        api = YouTubeTranscriptApi()
        transcript_list = api.fetch(video_id, languages=['en', 'en-US', 'en-GB'])

        transcript_lines = []
        for entry in transcript_list:
            start = entry.start
            text = entry.text.strip() if entry.text else ""
            if text:
                mins = int(start // 60)
                secs = int(start % 60)
                transcript_lines.append(f"[{mins}:{secs:02d}] {text}")

        if not transcript_lines:
            return {"error": "Transcript found but empty", "video_id": video_id}

        full_text = "\n".join(transcript_lines)

        return {
            "video_id": video_id,
            "language": "en",
            "line_count": len(transcript_lines),
            "transcript": full_text,
        }

    except Exception as e:
        error_msg = str(e)
        if "No transcripts" in error_msg or "TranscriptsDisabled" in error_msg:
            return {"error": "No English transcript available for this video", "video_id": video_id if video_id else url}
        return {"error": error_msg}


def tool_powerpoint_creator(args: Dict[str, Any]) -> Dict[str, Any]:
    """Generate a real .pptx file from a slide structure. Returns a download URL."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        return {"error": "python-pptx is not installed in this environment"}

    try:
        title = str(args.get("title", "Presentation"))
        subtitle = str(args.get("subtitle", ""))
        slides = args.get("slides", [])
        if not isinstance(slides, list):
            return {"error": "slides must be a list of {title, bullets, notes} objects"}

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Title slide
        title_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_layout)
        slide.shapes.title.text = title
        if subtitle and len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle

        # Content slides
        content_layout = prs.slide_layouts[1]
        for s in slides:
            if not isinstance(s, dict):
                continue
            slide = prs.slides.add_slide(content_layout)
            slide.shapes.title.text = str(s.get("title", ""))
            bullets = s.get("bullets", [])
            if bullets and len(slide.placeholders) > 1:
                tf = slide.placeholders[1].text_frame
                tf.text = ""
                first = True
                for b in bullets:
                    if first:
                        p = tf.paragraphs[0]
                        first = False
                    else:
                        p = tf.add_paragraph()
                    p.text = str(b)
                    p.font.size = Pt(20)
            notes = s.get("notes", "")
            if notes:
                slide.notes_slide.notes_text_frame.text = str(notes)

        # Save to storage/tool_outputs with a unique id
        output_dir = "/app/storage/tool_outputs"
        os.makedirs(output_dir, exist_ok=True)
        file_id = uuid.uuid4().hex[:12]
        safe_title = re.sub(r'[^a-zA-Z0-9_-]', '_', title)[:40] or "presentation"
        filename = f"{safe_title}.pptx"
        filepath = f"{output_dir}/{file_id}_{filename}"
        prs.save(filepath)

        return {
            "success": True,
            "download_url": f"/api/tools/download/{file_id}_{filename}",
            "filename": filename,
            "slide_count": len(slides) + 1,
        }
    except Exception as e:
        return {"error": str(e)}


REGISTRY: Dict[str, Dict[str, Any]] = {
    "calculator": {
        "schema": {
            "type": "object",
            "properties": {"expression": {"type": "string"}},
            "required": ["expression"]
        },
        "fn": tool_calculator
    },
    "youtube_transcript": {
        "schema": {
            "type": "object",
            "properties": {
                "url": {"type": "string", "description": "YouTube video URL (any format)"}
            },
            "required": ["url"]
        },
        "fn": tool_youtube_transcript
    },
    "powerpoint_creator": {
        "schema": {
            "type": "object",
            "properties": {
                "title": {"type": "string", "description": "Presentation title"},
                "subtitle": {"type": "string", "description": "Optional subtitle or author"},
                "slides": {
                    "type": "array",
                    "description": "Array of slide objects, each with title, bullets, and optional notes",
                    "items": {
                        "type": "object",
                        "properties": {
                            "title": {"type": "string"},
                            "bullets": {"type": "array", "items": {"type": "string"}},
                            "notes": {"type": "string"}
                        },
                        "required": ["title"]
                    }
                }
            },
            "required": ["title", "slides"]
        },
        "fn": tool_powerpoint_creator
    }
}

# Full Python runtime available to generated tools
_TOOL_IMPORTS = {
    "__builtins__": __builtins__,
    "json": json,
    "math": math,
    "re": re,
    "httpx": httpx,
    "os": os,
    "datetime": datetime,
    "timedelta": timedelta,
    "urlparse": urlparse,
    "parse_qs": parse_qs,
    "Dict": Dict,
    "Any": Any,
}


def make_dynamic_fn(code: str, fn_name: str):
    """Compile a generated tool function with full Python runtime access."""
    namespace: Dict[str, Any] = {**_TOOL_IMPORTS}
    compiled = compile(code, f'<tool:{fn_name}>', 'exec')
    exec(compiled, namespace)  # noqa: S102 — intentional dynamic tool loading
    if fn_name not in namespace:
        raise ValueError(f"Function '{fn_name}' not found in generated code")
    return namespace[fn_name]


def register_tool(name: str, schema: dict, code: str, fn_name: str):
    """Register a dynamically generated tool."""
    fn = make_dynamic_fn(code, fn_name)
    REGISTRY[name] = {"schema": schema, "fn": fn, "code": code, "generated": True}
    return REGISTRY[name]
